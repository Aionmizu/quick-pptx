"""Shared shape-creation helpers.

All native python-pptx calls live here. Per the architecture's implementation
patterns: no raw python-pptx API calls scattered across renderer/builder code —
they all go through this module so the python-pptx surface can be swapped
out cleanly if needed (e.g., for the hybrid Path C).
"""

from __future__ import annotations

from typing import TYPE_CHECKING

from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_AUTO_SIZE
from pptx.util import Emu, Pt

if TYPE_CHECKING:
    from pptx.presentation import Presentation
    from pptx.shapes.autoshape import Shape
    from pptx.slide import Slide

# Standard 16:9 slide dimensions (EMUs). 13.333 in × 7.5 in.
SLIDE_WIDTH_EMU: int = Emu(12192000)
SLIDE_HEIGHT_EMU: int = Emu(6858000)


def hex_to_rgb(hex_color: str) -> RGBColor:
    """Convert a `#RRGGBB` (or `RRGGBB`) hex string to a `RGBColor`."""
    h = hex_color.lstrip("#")
    if len(h) != 6:
        raise ValueError(f"Expected #RRGGBB hex color, got {hex_color!r}")
    r = int(h[0:2], 16)
    g = int(h[2:4], 16)
    b = int(h[4:6], 16)
    return RGBColor(r, g, b)


def add_solid_background(slide: Slide, hex_color: str) -> None:
    """Set the slide background to a solid hex color via fill."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = hex_to_rgb(hex_color)


def fit_title_font(
    text: str,
    base_pt: int,
    container_width_inches: float,
    container_height_inches: float,
    *,
    min_pt: int = 18,
) -> int:
    """Heuristically pick a title font size so the text wraps to ≤ 3 lines.

    The python-pptx `MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE` is unreliable across
    PowerPoint/Keynote/Google Slides — it relies on font metrics the renderer
    may not have. Pre-computing a sensible size from text length is more
    portable and produces predictable output everywhere the deck is opened.

    Heuristic:
      - Estimate chars-per-line at the base font size.
      - Compute lines needed for the text.
      - If >3 lines, scale the font down until ≤3 lines fit.
    """
    if not text:
        return base_pt
    n_chars = len(text)
    # Empirical: 1 inch of width fits ~ (200 / size_pt) characters in a sans-serif.
    # For Inter at 44pt → ~4.5 chars/inch; at 28pt → ~7 chars/inch.
    pt = base_pt
    while pt > min_pt:
        chars_per_inch = 200.0 / pt
        chars_per_line = max(1, int(container_width_inches * chars_per_inch))
        lines_needed = max(1, (n_chars + chars_per_line - 1) // chars_per_line)
        # Estimate vertical fit: line height ~ 1.2× font size in points;
        # 72 pt per inch.
        line_height_inches = (pt * 1.2) / 72.0
        vertical_fit = (lines_needed * line_height_inches) <= container_height_inches
        if lines_needed <= 3 and vertical_fit:
            return pt
        pt -= 4
    return min_pt


def add_text_box(
    slide: Slide,
    *,
    left_emu: int,
    top_emu: int,
    width_emu: int,
    height_emu: int,
    text: str,
    font_name: str,
    font_size_pt: int,
    font_color_hex: str,
    bold: bool = False,
    align_center: bool = False,
    auto_fit: bool = False,
) -> Shape:
    """Add a native, editable text box to a slide.

    Returns the underlying shape so callers can adjust further if needed.
    Native text — no rasterization (FR23).

    Args:
        auto_fit: when True, sets `MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE` so PowerPoint
            shrinks text on open if the renderer's pre-computed size is still
            too big. Use as a safety net alongside `fit_title_font`.
    """
    box = slide.shapes.add_textbox(left_emu, top_emu, width_emu, height_emu)
    tf = box.text_frame
    tf.word_wrap = True
    if auto_fit:
        tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    tf.text = text
    p = tf.paragraphs[0]
    if align_center:
        from pptx.enum.text import PP_ALIGN

        p.alignment = PP_ALIGN.CENTER
    for run in p.runs:
        run.font.name = font_name
        run.font.size = Pt(font_size_pt)
        run.font.bold = bold
        run.font.color.rgb = hex_to_rgb(font_color_hex)
    return box


def add_bullet_list(
    slide: Slide,
    *,
    left_emu: int,
    top_emu: int,
    width_emu: int,
    height_emu: int,
    bullets: list[str],
    font_name: str,
    font_size_pt: int,
    font_color_hex: str,
) -> Shape:
    """Add a textbox containing a vertical bullet list."""
    box = slide.shapes.add_textbox(left_emu, top_emu, width_emu, height_emu)
    tf = box.text_frame
    tf.word_wrap = True
    if not bullets:
        return box
    tf.text = bullets[0]
    for line in bullets[1:]:
        p = tf.add_paragraph()
        p.text = line
    for p in tf.paragraphs:
        for run in p.runs:
            run.font.name = font_name
            run.font.size = Pt(font_size_pt)
            run.font.color.rgb = hex_to_rgb(font_color_hex)
    return box


def add_color_block(
    slide: Slide,
    *,
    left_emu: int,
    top_emu: int,
    width_emu: int,
    height_emu: int,
    fill_hex: str,
) -> Shape:
    """Add a filled rectangle (used as accent block, divider band, etc.)."""
    from pptx.enum.shapes import MSO_SHAPE

    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left_emu, top_emu, width_emu, height_emu)
    shape.fill.solid()
    shape.fill.fore_color.rgb = hex_to_rgb(fill_hex)
    shape.line.fill.background()  # no border
    return shape


# OOXML core-property fields are capped at 255 characters by python-pptx.
# We keep the structured fields (read back by falsification) intact and
# truncate the rationale only.
_OOXML_CORE_PROP_LIMIT = 255
_RATIONALE_TRUNCATION_SUFFIX = "…[truncated]"


def write_metadata(presentation: Presentation, metadata: dict[str, str]) -> None:
    """Embed structural choices into the presentation's core+custom properties.

    Falsification readback later parses these out via `metadata.read_choices`.
    The rationale field is variable-length and may be truncated to fit
    within the 255-character core-property cap; structured fields are
    always preserved verbatim.
    """
    cp = presentation.core_properties
    # Keep title/subject minimal.
    cp.title = (metadata.get("style_name", "") or "")[:_OOXML_CORE_PROP_LIMIT]
    cp.subject = (metadata.get("layout_grid", "") or "")[:_OOXML_CORE_PROP_LIMIT]

    # Build the structured prefix first so we know how much room remains.
    structured_prefix = (
        f"layout_grid={metadata.get('layout_grid', '')}; "
        f"section_structure={metadata.get('section_structure', '')}; "
        f"hierarchy_pattern={metadata.get('hierarchy_pattern', '')}; "
        f"content_density={metadata.get('content_density', '')}; "
        f"style_name={metadata.get('style_name', '')}; "
        f"rationale="
    )
    rationale = metadata.get("rationale", "") or ""

    # Compute available space for rationale, leaving room for truncation suffix.
    available_for_rationale = _OOXML_CORE_PROP_LIMIT - len(structured_prefix)
    if len(rationale) <= available_for_rationale:
        cp.comments = structured_prefix + rationale
    else:
        # Reserve suffix length within the budget so the marker fits.
        max_rationale_len = max(0, available_for_rationale - len(_RATIONALE_TRUNCATION_SUFFIX))
        truncated = rationale[:max_rationale_len].rstrip() + _RATIONALE_TRUNCATION_SUFFIX
        cp.comments = (structured_prefix + truncated)[:_OOXML_CORE_PROP_LIMIT]


__all__ = [
    "SLIDE_HEIGHT_EMU",
    "SLIDE_WIDTH_EMU",
    "add_bullet_list",
    "add_color_block",
    "add_solid_background",
    "add_text_box",
    "fit_title_font",
    "hex_to_rgb",
    "write_metadata",
]
