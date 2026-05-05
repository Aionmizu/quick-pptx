"""Slide renderer — emits native python-pptx output from `SlidePlan`s + tokens.

Supports the four core layout grids: single-column, two-up, asymmetric, bento,
plus a dedicated full-bleed title slide for section dividers.

All shape creation goes through `ia_pptx.core.shapes` so the python-pptx
surface remains isolated.

Layout philosophy (post-visual-quality fix):
- Titles use `fit_title_font` to scale font down for long French/English titles
  that would otherwise wrap one-word-per-line.
- Asymmetric layout uses a 40/60 split (was 30/70) so titles have more room.
- Section dividers are full-bleed accent slides — generous typography, not
  narrow column-on-strip.
- Auto-fit (`auto_fit=True` on title boxes) catches cases the heuristic missed.
"""

from __future__ import annotations

import logging
from pathlib import Path

from pptx import Presentation
from pptx.util import Emu

from ia_pptx.core.exceptions import RenderFailed
from ia_pptx.core.shapes import (
    SLIDE_HEIGHT_EMU,
    SLIDE_WIDTH_EMU,
    add_bullet_list,
    add_color_block,
    add_solid_background,
    add_text_box,
    fit_title_font,
    write_metadata,
)
from ia_pptx.core.types import (
    DesignTokens,
    LayoutGrid,
    SlidePlan,
    StructuralChoices,
)

logger = logging.getLogger(__name__)

_INCH = 914400  # EMU per inch


def _typography_sizes(content_density: str) -> tuple[int, int, int]:
    """Title / body / accent font sizes in points, scaled by density."""
    if content_density == "minimal":
        return 56, 22, 18
    if content_density == "dense":
        return 36, 16, 14
    return 44, 18, 16


def _render_title_slide(
    slide,
    plan: SlidePlan,
    tokens: DesignTokens,
) -> None:
    """Render a section-divider / opening title slide.

    Full-bleed background with the accent color as a strong horizontal band,
    centered title in display weight, optional subtitle below.
    """
    title_pt_base, body_pt, _ = _typography_sizes(tokens.content_density.value)
    add_solid_background(slide, tokens.palette.background_hex)

    # A wide accent band gives the section divider visual weight without the
    # narrow-strip-with-tiny-column look of the previous version.
    add_color_block(
        slide,
        left_emu=Emu(int(0.0 * _INCH)),
        top_emu=Emu(int(2.6 * _INCH)),
        width_emu=Emu(int(13.333 * _INCH)),
        height_emu=Emu(int(0.18 * _INCH)),
        fill_hex=tokens.palette.accent_hex,
    )

    # Title — full slide width, generously sized; auto-fit as safety net.
    title_w = 12.0
    title_h = 2.2
    title_pt = fit_title_font(plan.title, title_pt_base, title_w, title_h, min_pt=24)
    add_text_box(
        slide,
        left_emu=Emu(int(0.7 * _INCH)),
        top_emu=Emu(int(0.5 * _INCH)),
        width_emu=Emu(int(title_w * _INCH)),
        height_emu=Emu(int(title_h * _INCH)),
        text=plan.title,
        font_name=tokens.typography.heading_font,
        font_size_pt=title_pt,
        font_color_hex=tokens.palette.text_hex,
        bold=True,
        auto_fit=True,
    )
    if plan.body:
        add_text_box(
            slide,
            left_emu=Emu(int(0.7 * _INCH)),
            top_emu=Emu(int(2.95 * _INCH)),
            width_emu=Emu(int(12.0 * _INCH)),
            height_emu=Emu(int(3.0 * _INCH)),
            text="\n".join(plan.body),
            font_name=tokens.typography.body_font,
            font_size_pt=body_pt + 4,
            font_color_hex=tokens.palette.secondary_hex,
        )


def _render_single_column(slide, plan: SlidePlan, tokens: DesignTokens) -> None:
    title_pt_base, body_pt, _ = _typography_sizes(tokens.content_density.value)
    add_solid_background(slide, tokens.palette.background_hex)

    title_w = 12.0
    title_h = 1.6
    title_pt = fit_title_font(plan.title, title_pt_base - 8, title_w, title_h, min_pt=22)
    add_text_box(
        slide,
        left_emu=Emu(int(0.7 * _INCH)),
        top_emu=Emu(int(0.6 * _INCH)),
        width_emu=Emu(int(title_w * _INCH)),
        height_emu=Emu(int(title_h * _INCH)),
        text=plan.title,
        font_name=tokens.typography.heading_font,
        font_size_pt=title_pt,
        font_color_hex=tokens.palette.text_hex,
        bold=True,
        auto_fit=True,
    )
    # Thin accent rule under the title for visual structure.
    add_color_block(
        slide,
        left_emu=Emu(int(0.7 * _INCH)),
        top_emu=Emu(int(2.3 * _INCH)),
        width_emu=Emu(int(2.0 * _INCH)),
        height_emu=Emu(int(0.06 * _INCH)),
        fill_hex=tokens.palette.accent_hex,
    )
    add_bullet_list(
        slide,
        left_emu=Emu(int(0.7 * _INCH)),
        top_emu=Emu(int(2.6 * _INCH)),
        width_emu=Emu(int(12 * _INCH)),
        height_emu=Emu(int(4.5 * _INCH)),
        bullets=list(plan.body),
        font_name=tokens.typography.body_font,
        font_size_pt=body_pt,
        font_color_hex=tokens.palette.text_hex,
    )


def _render_two_up(slide, plan: SlidePlan, tokens: DesignTokens) -> None:
    title_pt_base, body_pt, _ = _typography_sizes(tokens.content_density.value)
    add_solid_background(slide, tokens.palette.background_hex)

    title_w = 12.0
    title_h = 1.4
    title_pt = fit_title_font(plan.title, title_pt_base - 10, title_w, title_h, min_pt=22)
    add_text_box(
        slide,
        left_emu=Emu(int(0.7 * _INCH)),
        top_emu=Emu(int(0.6 * _INCH)),
        width_emu=Emu(int(title_w * _INCH)),
        height_emu=Emu(int(title_h * _INCH)),
        text=plan.title,
        font_name=tokens.typography.heading_font,
        font_size_pt=title_pt,
        font_color_hex=tokens.palette.text_hex,
        bold=True,
        auto_fit=True,
    )
    # Thin accent rule under the title.
    add_color_block(
        slide,
        left_emu=Emu(int(0.7 * _INCH)),
        top_emu=Emu(int(2.1 * _INCH)),
        width_emu=Emu(int(1.6 * _INCH)),
        height_emu=Emu(int(0.06 * _INCH)),
        fill_hex=tokens.palette.accent_hex,
    )
    # Two side-by-side columns, splitting body roughly in half.
    half = max(1, len(plan.body) // 2)
    left_bullets = list(plan.body[:half])
    right_bullets = list(plan.body[half:])
    add_bullet_list(
        slide,
        left_emu=Emu(int(0.7 * _INCH)),
        top_emu=Emu(int(2.4 * _INCH)),
        width_emu=Emu(int(5.8 * _INCH)),
        height_emu=Emu(int(4.7 * _INCH)),
        bullets=left_bullets,
        font_name=tokens.typography.body_font,
        font_size_pt=body_pt,
        font_color_hex=tokens.palette.text_hex,
    )
    add_bullet_list(
        slide,
        left_emu=Emu(int(7.0 * _INCH)),
        top_emu=Emu(int(2.4 * _INCH)),
        width_emu=Emu(int(5.8 * _INCH)),
        height_emu=Emu(int(4.7 * _INCH)),
        bullets=right_bullets,
        font_name=tokens.typography.body_font,
        font_size_pt=body_pt,
        font_color_hex=tokens.palette.text_hex,
    )


def _render_asymmetric(slide, plan: SlidePlan, tokens: DesignTokens) -> None:
    """Asymmetric: 40/60 split — colored panel on the left, content on the right.

    Title sits in the colored panel with auto-scaled font; body on the right.
    """
    title_pt_base, body_pt, _ = _typography_sizes(tokens.content_density.value)
    add_solid_background(slide, tokens.palette.background_hex)

    panel_width = 5.3  # was 4.5 — give titles room to breathe
    add_color_block(
        slide,
        left_emu=Emu(int(0.0 * _INCH)),
        top_emu=Emu(int(0.0 * _INCH)),
        width_emu=Emu(int(panel_width * _INCH)),
        height_emu=Emu(int(7.5 * _INCH)),
        fill_hex=tokens.palette.primary_hex,
    )

    # Title inside the panel; widen to ~80% of panel width.
    title_w = panel_width - 1.0
    title_h = 5.2
    title_pt = fit_title_font(plan.title, title_pt_base, title_w, title_h, min_pt=22)
    add_text_box(
        slide,
        left_emu=Emu(int(0.5 * _INCH)),
        top_emu=Emu(int(0.8 * _INCH)),
        width_emu=Emu(int(title_w * _INCH)),
        height_emu=Emu(int(title_h * _INCH)),
        text=plan.title,
        font_name=tokens.typography.heading_font,
        font_size_pt=title_pt,
        font_color_hex=tokens.palette.background_hex,
        bold=True,
        auto_fit=True,
    )
    # Body on the right column.
    add_bullet_list(
        slide,
        left_emu=Emu(int((panel_width + 0.4) * _INCH)),
        top_emu=Emu(int(0.8 * _INCH)),
        width_emu=Emu(int((13.0 - panel_width - 0.6) * _INCH)),
        height_emu=Emu(int(6.0 * _INCH)),
        bullets=list(plan.body),
        font_name=tokens.typography.body_font,
        font_size_pt=body_pt,
        font_color_hex=tokens.palette.text_hex,
    )


def _render_bento(slide, plan: SlidePlan, tokens: DesignTokens) -> None:
    title_pt_base, body_pt, _ = _typography_sizes(tokens.content_density.value)
    add_solid_background(slide, tokens.palette.background_hex)

    title_w = 12.0
    title_h = 1.0
    title_pt = fit_title_font(plan.title, title_pt_base - 12, title_w, title_h, min_pt=20)
    add_text_box(
        slide,
        left_emu=Emu(int(0.6 * _INCH)),
        top_emu=Emu(int(0.4 * _INCH)),
        width_emu=Emu(int(title_w * _INCH)),
        height_emu=Emu(int(title_h * _INCH)),
        text=plan.title,
        font_name=tokens.typography.heading_font,
        font_size_pt=title_pt,
        font_color_hex=tokens.palette.text_hex,
        bold=True,
        auto_fit=True,
    )
    # Bento composition: 3 tiles of varying sizes
    bullets = list(plan.body) or ["", "", ""]
    while len(bullets) < 3:
        bullets.append("")
    # Big tile (top-left)
    add_color_block(
        slide,
        left_emu=Emu(int(0.6 * _INCH)),
        top_emu=Emu(int(1.6 * _INCH)),
        width_emu=Emu(int(7.0 * _INCH)),
        height_emu=Emu(int(3.5 * _INCH)),
        fill_hex=tokens.palette.accent_hex,
    )
    add_text_box(
        slide,
        left_emu=Emu(int(0.9 * _INCH)),
        top_emu=Emu(int(2.0 * _INCH)),
        width_emu=Emu(int(6.4 * _INCH)),
        height_emu=Emu(int(3.0 * _INCH)),
        text=bullets[0],
        font_name=tokens.typography.body_font,
        font_size_pt=body_pt + 4,
        font_color_hex=tokens.palette.background_hex,
        bold=True,
        auto_fit=True,
    )
    # Top-right tile
    add_color_block(
        slide,
        left_emu=Emu(int(7.9 * _INCH)),
        top_emu=Emu(int(1.6 * _INCH)),
        width_emu=Emu(int(4.8 * _INCH)),
        height_emu=Emu(int(1.6 * _INCH)),
        fill_hex=tokens.palette.primary_hex,
    )
    add_text_box(
        slide,
        left_emu=Emu(int(8.1 * _INCH)),
        top_emu=Emu(int(1.8 * _INCH)),
        width_emu=Emu(int(4.4 * _INCH)),
        height_emu=Emu(int(1.2 * _INCH)),
        text=bullets[1],
        font_name=tokens.typography.body_font,
        font_size_pt=body_pt,
        font_color_hex=tokens.palette.background_hex,
        auto_fit=True,
    )
    # Bottom-right tile
    add_color_block(
        slide,
        left_emu=Emu(int(7.9 * _INCH)),
        top_emu=Emu(int(3.4 * _INCH)),
        width_emu=Emu(int(4.8 * _INCH)),
        height_emu=Emu(int(1.7 * _INCH)),
        fill_hex=tokens.palette.secondary_hex,
    )
    add_text_box(
        slide,
        left_emu=Emu(int(8.1 * _INCH)),
        top_emu=Emu(int(3.6 * _INCH)),
        width_emu=Emu(int(4.4 * _INCH)),
        height_emu=Emu(int(1.3 * _INCH)),
        text=bullets[2],
        font_name=tokens.typography.body_font,
        font_size_pt=body_pt,
        font_color_hex=tokens.palette.background_hex,
        auto_fit=True,
    )
    # Wide bottom tile if there are more bullets
    if len(bullets) > 3 and bullets[3]:
        add_text_box(
            slide,
            left_emu=Emu(int(0.6 * _INCH)),
            top_emu=Emu(int(5.4 * _INCH)),
            width_emu=Emu(int(12 * _INCH)),
            height_emu=Emu(int(1.6 * _INCH)),
            text="\n".join(bullets[3:]),
            font_name=tokens.typography.body_font,
            font_size_pt=body_pt,
            font_color_hex=tokens.palette.text_hex,
        )


_LAYOUT_DISPATCH: dict[
    LayoutGrid,
    object,  # callable; typed loosely to avoid forward-ref noise
] = {
    LayoutGrid.SINGLE_COLUMN: _render_single_column,
    LayoutGrid.TWO_UP: _render_two_up,
    LayoutGrid.ASYMMETRIC: _render_asymmetric,
    LayoutGrid.BENTO: _render_bento,
}


def render_deck(
    slides: list[SlidePlan],
    tokens: DesignTokens,
    output_path: Path,
    structural_metadata: dict[str, str] | None = None,
) -> Path:
    """Render the slide plans to a `.pptx` at `output_path`.

    Returns the absolute output path. Raises `RenderFailed` on errors.
    """
    if not slides:
        raise RenderFailed("Cannot render deck with zero slides")
    try:
        prs = Presentation()
        # Set slide size to 16:9 widescreen
        prs.slide_width = SLIDE_WIDTH_EMU
        prs.slide_height = SLIDE_HEIGHT_EMU
        blank_layout = prs.slide_layouts[6]  # blank layout

        for plan in slides:
            slide = prs.slides.add_slide(blank_layout)
            if plan.is_section_divider:
                _render_title_slide(slide, plan, tokens)
                continue
            handler = _LAYOUT_DISPATCH.get(plan.layout_grid, _render_single_column)
            handler(slide, plan, tokens)  # type: ignore[operator]
            if plan.notes:
                slide.notes_slide.notes_text_frame.text = plan.notes

        if structural_metadata is not None:
            write_metadata(prs, structural_metadata)

        output_path.parent.mkdir(parents=True, exist_ok=True)
        prs.save(str(output_path))
        return output_path.resolve()
    except RenderFailed:
        raise
    except Exception as exc:
        raise RenderFailed(f"Render failed: {exc}") from exc


def tokens_from_choices(choices: StructuralChoices) -> DesignTokens:
    """Project structural choices into render-time tokens."""
    return DesignTokens(
        palette=choices.style.palette,
        typography=choices.style.typography,
        layout_grid=choices.layout_grid,
        hierarchy_pattern=choices.hierarchy_pattern,
        content_density=choices.content_density,
    )


__all__ = [
    "render_deck",
    "tokens_from_choices",
]
