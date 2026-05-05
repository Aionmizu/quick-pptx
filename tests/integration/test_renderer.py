"""Integration tests for the renderer.

Each layout grid is rendered to a real `.pptx`, then re-opened and inspected
to verify text content is editable (native, not rasterized).
"""

from __future__ import annotations

from pathlib import Path

import pytest
from pptx import Presentation

from ia_pptx.core.renderer import render_deck, tokens_from_choices
from ia_pptx.core.types import (
    ContentDensity,
    HierarchyPattern,
    LayoutGrid,
    SectionStructure,
    SlidePlan,
    StructuralChoices,
)
from ia_pptx.design import get_design_library


def _make_choices(layout: LayoutGrid) -> StructuralChoices:
    lib = get_design_library()
    style = lib.sample_style(seed=7)
    return StructuralChoices(
        layout_grid=layout,
        section_structure=SectionStructure.LINEAR,
        hierarchy_pattern=HierarchyPattern.TYPE_LED,
        content_density=ContentDensity.MEDIUM,
        style=style,
        rationale="test",
    )


def _make_plans(layout: LayoutGrid) -> list[SlidePlan]:
    return [
        SlidePlan(
            layout_grid=layout,
            is_section_divider=True,
            title="Sample Deck",
            body=("Subtitle line",),
        ),
        SlidePlan(
            layout_grid=layout,
            is_section_divider=False,
            title="Why this matters",
            body=(
                "First key point goes here",
                "Second key point complements the first",
                "Third key point closes the section",
            ),
        ),
        SlidePlan(
            layout_grid=layout,
            is_section_divider=False,
            title="Key numbers",
            body=("88% of users do this", "12 minutes saved per task"),
        ),
    ]


@pytest.mark.parametrize(
    "layout",
    [
        LayoutGrid.SINGLE_COLUMN,
        LayoutGrid.TWO_UP,
        LayoutGrid.ASYMMETRIC,
        LayoutGrid.BENTO,
    ],
)
def test_renders_each_layout_grid(layout: LayoutGrid, tmp_path: Path) -> None:
    choices = _make_choices(layout)
    plans = _make_plans(layout)
    tokens = tokens_from_choices(choices)
    out = tmp_path / f"{layout.value}.pptx"
    written = render_deck(plans, tokens, out, structural_metadata=choices.to_metadata_dict())
    assert written.exists()
    assert written.stat().st_size > 0

    # Re-open and verify text is native + present.
    prs = Presentation(str(written))
    assert len(prs.slides) == 3
    # Pull all text content out and check the first slide's title is there.
    slide_texts: list[str] = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                slide_texts.append(shape.text_frame.text)
    blob = "\n".join(slide_texts)
    assert "Sample Deck" in blob
    assert "Why this matters" in blob
    assert "Key numbers" in blob


def test_metadata_is_recorded_in_core_properties(tmp_path: Path) -> None:
    choices = _make_choices(LayoutGrid.BENTO)
    plans = _make_plans(LayoutGrid.BENTO)
    tokens = tokens_from_choices(choices)
    out = tmp_path / "metadata.pptx"
    render_deck(plans, tokens, out, structural_metadata=choices.to_metadata_dict())
    prs = Presentation(str(out))
    cp = prs.core_properties
    # Title carries style name; comments carry the full structural choice block.
    assert cp.title == choices.style.name
    assert "layout_grid=bento" in cp.comments
    assert "section_structure=linear" in cp.comments
